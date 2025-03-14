import h2ogpte
from h2ogpte import H2OGPTE
from h2o_wave import main, app, Q, ui, on, run_on, data
import hashlib
import os
import json
from pptx import Presentation
from pptx.util import Pt
from ftfy import fix_text


@app('/')
async def serve(q: Q):
    # First time a browser comes to the app
    
    if not q.client.initialized:
        await init(q)
        q.client.initialized = True

    q.page['meta'].notification_bar = None

    # Other browser interactions
    await run_on(q)
    await q.page.save()

async def init(q: Q) -> None:
    q.client.cards = set()
    q.app.icon, = await q.site.upload(['./static/icon.png'])

    q.page['meta'] = ui.meta_card(
        box='',
        title='Automatic Slideshow Generator',
        theme='dark',
        script=heap_analytics(
            userid=q.auth.subject,
            event_properties=f"{{"
                             f"version: '0.0.1', "
                             f"product: 'Slideshow_Generator'"
                             f"}}",
        ),
        layouts=[
            ui.layout(
                breakpoint='xs',
                min_height='100vh',
                max_width='1200px',
                zones=[
                    ui.zone('header'),
                    ui.zone('content', size='1', zones=[
                        ui.zone('vertical', size='1',),
                        ui.zone('horizontal', direction=ui.ZoneDirection.ROW),
                        ui.zone('grid', direction=ui.ZoneDirection.ROW, wrap='stretch', justify='center')
                    ]),
                    ui.zone(name='footer'),
                ]
            )
        ]
    )

    q.page['header'] = ui.header_card(
        box='header',
        title='Automatic Slideshow Generator',
        subtitle="Creates an AI-generated powerpoint using h2oGPTe!",
        image=q.app.icon
    )
    q.page['footer'] = ui.footer_card(
        box='footer',
        caption='Made with 💛 using [H2O Wave](https://wave.h2o.ai).'
    )

    await home(q)

@on()
async def home(q: Q):
    
    clear_cards(q)

    items_intro = [ui.text_xl('**Welcome! Get Ready to Create Your Automated Powerpoint!**'),]

    items_below = await get_input_items(q)
    add_card(q,'intro',ui.form_card(box='vertical', items=[
        ui.inline(items= items_intro,justify='center'),
        ui.separator(),
        ui.inline(items=[ui.image(title='gif',path='https://i.postimg.cc/W4jVNfhm/giphy.gif',width = '700px'),],justify='center'),
        ui.separator(),
        ui.inline(items_below,justify='center'),
        ui.inline(items=[
            ui.button(name='save', label = 'Save', primary = True),
            ui.button(name='refresh', label = 'Refresh')],justify='center')]))
    # if q.client.demo_mode:
    #     await call_demo_actions(q)
        
@on()
async def refresh(q: Q):
    await home(q)


async def get_input_items(q: Q):
   
    items = [ui.textbox(name='url',width = '300px', label='Choose Your Topic', placeholder='i.e. dogs, earth, etc'),
             ui.dropdown(
                    name="slides",
                    # value=None if q.client.climate_zone is None else q.client.climate_zone,
                    # trigger=True,
                    choices=[
                        ui.choice(name="1", label="1"),
                        ui.choice(name="2", label="2"),
                        ui.choice(name="3", label="3"),
                        ui.choice(name="4", label="4"),
                        ui.choice(name="5", label="5"),
                        ui.choice(name="6", label="6"),
                        ui.choice(name="7", label="7"),
                    ],
                    placeholder="Pick me!",
                    label="Number of Slides",
                    width="300px"),
            ui.dropdown(
                    name="bullets",
                    # value=None if q.client.climate_zone is None else q.client.climate_zone,
                    # trigger=True,
                    choices=[
                        ui.choice(name="1", label="1"),
                        ui.choice(name="2", label="2"),
                        ui.choice(name="3", label="3"),
                        ui.choice(name="4", label="4"),
                        ui.choice(name="5", label="5"),
                    ],
                    placeholder="Pick me!",
                    label="Number of Bullet Points Per Slide",
                    width="300px"),
            ]

    return items

@on()
async def save(q: Q):
    """Answer the question"""
    #Gets youtube link, and uses API to get the transcript
    q.topic = q.args.url
    q.slides = q.args.slides
    q.bullets = q.args.bullets
    q.page["intro"].data += [q.topic, True]
    q.page["intro"].data += [q.slides, True]
    q.page["intro"].data += [q.bullets, True]
    await q.page.save()
    await long_process_dialog(q)

    
async def long_process_dialog(q):
    q.app.dialog_gif, = await q.site.upload(['./static/load.gif'])
    items = [
        ui.image(title="", path=q.app.dialog_gif, width="550px"),
    ]
    q.page["meta"].dialog = ui.dialog(
        title='Generating a Powerpoint About Your Topic...',
        items=items,
        blocking=True
    )
    await q.page.save()
    q.page["meta"].dialog = None
    await get_json_format(q)

async def get_json_format(q:Q):
    # Start the chat session for this one document within the same collection
    h2ogpte_url = os.getenv("h2ogpte_url")
    q.client.address = h2ogpte_url
    h2ogpte_key = os.getenv("h2ogpte_key")
    q.client.h2ogpte_key = h2ogpte_key
    client = H2OGPTE(address=h2ogpte_url, api_key=h2ogpte_key)
    collection_id = client.create_collection(name='slideshow_creator', 
                                                 description="Collection for document: slideshow_creator"
                                                 )
    chat_session_id = client.create_chat_session(collection_id=collection_id)

    output1 = {
    "Slide1": {
    "title": "Introduction",
    "bullet1": "Dogs are domesticated mammals and are a subspecies of the wolf..."
    }}
    output2 = str(output1)
    # Engineer the Prompt 
    t_system_prompt = "You are h2oGPTe, an expert question-answering AI system created by H2O.ai that performs like GPT-4 by OpenAI.\n"
    t_llm = "mistralai/Mixtral-8x7B-Instruct-v0.1"

    with client.connect(chat_session_id) as session:
        secexcl = session.query(
            system_prompt=t_system_prompt,
            pre_prompt_query=None,
            prompt_query=None,
            message= "create a powerpoint template with " + q.slides + " slides of " + q.bullets + " bullet points about the topic of " + q.topic + " in a json format like this: " + output2 + "Make sure the output is a json file or else I will lose my job. Do NOT explain yourself or rephrase the question. The output should just be in a JSON format or else I will lose my job.",
            llm=t_llm,
            llm_args={"max_new_tokens":3072},
            rag_config={"rag_type": "llm_only"},
            timeout=2000,
        
        ).content
        # Only returns for one section here
        print(secexcl)
        fixed_text = fix_text(secexcl)

        with open('slideshow.txt', "w") as f:
            f.write(fixed_text)
        with open('slideshow.txt') as f:
            output = json.load(f)
            print(output)

    num_slides = int(q.slides)
    num_bullets = int(q.bullets)

    prs = Presentation("./ppt_template.pptx")
    bullet_slide_layout = prs.slide_layouts[1]

    for i in range(num_slides):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        string1 = i + 1
        overall_slide = "Slide"+str(string1)
        title_shape.text = output[overall_slide]['title']
        tf = body_shape.text_frame
        for j in range(num_bullets):
            p = tf.add_paragraph()
            run = p.add_run()
            string2 = j + 1
            overall_bullet = "bullet"+str(string2)
            run.text = output[overall_slide][overall_bullet]
            font = run.font
            # font.name = 'Calibri'
            font.size = Pt(14)

    prs.save('test.pptx')

    final_text = "**Your AI-Generated Presentation Is Ready!**"
    items = [
        ui.inline([ui.text_xl(final_text)], justify='center'),
        ui.separator()
        ]
    download_path, = await q.site.upload(['./test.pptx'])
    items += [
            ui.inline(items=[ui.image(title='gif',path='https://i.postimg.cc/ZRVs5FQg/friends.webp',width = '600px'),],justify='center'),
            ui.separator(),
            ui.inline([
                 ui.link(name='download_pdf', label='View AI-Generated Powerpoint', path=download_path, download=True,  button=True),
            ], justify='center'),
            ]
    add_card(q, 'intro', ui.form_card(box='vertical',items=items))

def heap_analytics(userid, event_properties=None) -> ui.inline_script:

    if "HEAP_ID" not in os.environ:
        return
        
    heap_id = os.getenv("HEAP_ID")
    script = f"""
window.heap=window.heap||[],heap.load=function(e,t){{window.heap.appid=e,window.heap.config=t=t||{{}};var r=document.createElement("script");r.type="text/javascript",r.async=!0,r.src="https://cdn.heapanalytics.com/js/heap-"+e+".js";var a=document.getElementsByTagName("script")[0];a.parentNode.insertBefore(r,a);for(var n=function(e){{return function(){{heap.push([e].concat(Array.prototype.slice.call(arguments,0)))}}}},p=["addEventProperties","addUserProperties","clearEventProperties","identify","resetIdentity","removeEventProperty","setEventProperties","track","unsetEventProperty"],o=0;o<p.length;o++)heap[p[o]]=n(p[o])}};
heap.load("{heap_id}"); 
    """

    if userid is not None:  # is OIDC Enabled? we do not want to identify all non-logged in users as "none"
        identity = hashlib.sha256(userid.encode()).hexdigest()
        script += f"heap.identify('{identity}');"

    if event_properties is not None:
        script += f"heap.addEventProperties({event_properties})"

    return ui.inline_script(content=script)

def add_card(q, name, card) -> None:
    q.client.cards.add(name)
    q.page[name] = card

def clear_cards(q, ignore=[]) -> None:
    for name in q.client.cards.copy():
        if name not in ignore:
            del q.page[name]
            q.client.cards.remove(name)